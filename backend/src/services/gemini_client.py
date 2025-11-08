"""Gemini AI client for multimodal flashcard generation.

Supports:
- Text extraction from PDFs
- Image analysis (JPG, PNG)
- PowerPoint text extraction
- Direct text input

Requires:
    - GOOGLE_API_KEY environment variable
    - pip install google-generativeai PyPDF2 python-pptx Pillow
"""
import os
from typing import List, Dict
from pathlib import Path
import google.generativeai as genai
from dotenv import load_dotenv

load_dotenv()

# Configure Gemini API
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if GOOGLE_API_KEY:
    genai.configure(api_key=GOOGLE_API_KEY)

DEFAULT_MODEL = "gemini-1.5-flash"  # or "gemini-1.5-pro" for better quality


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file."""
    try:
        import PyPDF2
        text_parts = []
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        return "\n\n".join(text_parts)
    except Exception as e:
        raise RuntimeError(f"Failed to extract text from PDF: {e}")


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint file."""
    try:
        from pptx import Presentation
        text_parts = []
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        return "\n\n".join(text_parts)
    except Exception as e:
        raise RuntimeError(f"Failed to extract text from PPTX: {e}")


def process_image_with_gemini(file_path: str, prompt: str) -> str:
    """Send image to Gemini vision model and get response."""
    try:
        from PIL import Image
        img = Image.open(file_path)
        model = genai.GenerativeModel(DEFAULT_MODEL)
        response = model.generate_content([prompt, img])
        return response.text
    except Exception as e:
        raise RuntimeError(f"Failed to process image with Gemini: {e}")


def generate_flashcards_from_text(text: str, count: int = 10, model_name: str = DEFAULT_MODEL) -> List[Dict]:
    """Generate flashcards from text using Gemini.

    Args:
        text: Source text to generate flashcards from
        count: Number of flashcards to generate
        model_name: Gemini model to use

    Returns:
        List of dicts with 'question' and 'answer' keys
    """
    if not GOOGLE_API_KEY:
        raise RuntimeError("GOOGLE_API_KEY not set in environment")

    # Construct prompt
    prompt = f"""You are an expert educator creating study flashcards.

Given the following content, create exactly {count} high-quality flashcards.

Requirements:
- Each flashcard must have a clear question and concise answer
- Focus on key concepts, definitions, and important facts
- Make questions specific and answerable
- Return ONLY a valid JSON array with this exact format: [{{"question": "...", "answer": "..."}}, ...]
- Do not include any markdown formatting, code blocks, or additional text

Content:
{text[:8000]}

Generate {count} flashcards as a JSON array:"""

    try:
        model = genai.GenerativeModel(model_name)
        response = model.generate_content(prompt)

        # Extract JSON from response
        import json
        import re

        response_text = response.text.strip()

        # Try to find JSON array in response
        # Remove markdown code blocks if present
        response_text = re.sub(r'```json\s*', '', response_text)
        response_text = re.sub(r'```\s*$', '', response_text)

        # Find JSON array
        match = re.search(r'\[.*\]', response_text, re.DOTALL)
        if match:
            json_str = match.group(0)
            flashcards = json.loads(json_str)

            # Validate and normalize
            validated = []
            for card in flashcards:
                if isinstance(card, dict) and 'question' in card and 'answer' in card:
                    validated.append({
                        'question': str(card['question']).strip(),
                        'answer': str(card['answer']).strip()
                    })

            return validated[:count]
        else:
            # Fallback: try parsing the whole response
            flashcards = json.loads(response_text)
            validated = []
            for card in flashcards:
                if isinstance(card, dict) and 'question' in card and 'answer' in card:
                    validated.append({
                        'question': str(card['question']).strip(),
                        'answer': str(card['answer']).strip()
                    })
            return validated[:count]

    except Exception as e:
        raise RuntimeError(f"Failed to generate flashcards with Gemini: {e}")


def generate_flashcards_from_file(file_path: str, mime_type: str, count: int = 10) -> List[Dict]:
    """Generate flashcards from uploaded file.

    Automatically detects file type and uses appropriate processing.

    Args:
        file_path: Path to uploaded file
        mime_type: MIME type of file
        count: Number of flashcards to generate

    Returns:
        List of flashcard dicts
    """
    file_ext = Path(file_path).suffix.lower()

    # Handle PDFs
    if mime_type == 'application/pdf' or file_ext == '.pdf':
        text = extract_text_from_pdf(file_path)
        return generate_flashcards_from_text(text, count)

    # Handle PowerPoint
    elif mime_type in ['application/vnd.ms-powerpoint',
                       'application/vnd.openxmlformats-officedocument.presentationml.presentation'] \
         or file_ext in ['.ppt', '.pptx']:
        text = extract_text_from_pptx(file_path)
        return generate_flashcards_from_text(text, count)

    # Handle images - use vision model
    elif mime_type.startswith('image/') or file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.webp']:
        prompt = f"""Analyze this image and extract all text, concepts, and information.
Then create {count} educational flashcards based on the content.
Return ONLY a JSON array: [{{"question": "...", "answer": "..."}}, ...]"""

        response = process_image_with_gemini(file_path, prompt)

        # Parse JSON from response
        import json
        import re
        match = re.search(r'\[.*\]', response, re.DOTALL)
        if match:
            return json.loads(match.group(0))[:count]
        else:
            # If Gemini didn't return structured data, extract text and regenerate
            text_prompt = "Extract all text and information from this image as plain text."
            extracted_text = process_image_with_gemini(file_path, text_prompt)
            return generate_flashcards_from_text(extracted_text, count)

    # Handle text files
    elif mime_type.startswith('text/') or file_ext in ['.txt', '.md']:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        return generate_flashcards_from_text(text, count)

    else:
        raise ValueError(f"Unsupported file type: {mime_type}")
